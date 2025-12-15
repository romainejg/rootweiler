# imaging_tools.py

import io
import os
import zipfile
from typing import List, Tuple

import streamlit as st
from PIL import Image, ImageDraw
import numpy as np

# For document image extraction
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import openpyxl

# For canvas drawing
try:
    from streamlit_drawable_canvas import st_canvas
    CANVAS_AVAILABLE = True
except ImportError:
    CANVAS_AVAILABLE = False


# ----------------- Helper functions for Object Crop tool ----------------- #

def crop_square_png(image: Image.Image, bbox: Tuple[int, int, int, int], pad_factor: float = 1.1) -> Image.Image:
    """
    Crop a square region from an image centered on the given bounding box.
    
    Args:
        image: Input PIL Image
        bbox: Bounding box as (x0, y0, x1, y1) - top-left and bottom-right corners
        pad_factor: Padding factor to expand the square (default 1.1 = 110% of bbox size)
    
    Returns:
        Square PIL Image in RGBA mode with transparent padding if needed
    """
    x0, y0, x1, y1 = bbox
    
    # Normalize coordinates
    x0, x1 = min(x0, x1), max(x0, x1)
    y0, y1 = min(y0, y1), max(y0, y1)
    
    # Calculate bbox dimensions and center
    bbox_width = x1 - x0
    bbox_height = y1 - y0
    center_x = (x0 + x1) / 2
    center_y = (y0 + y1) / 2
    
    # Calculate square side length (max dimension * padding factor)
    square_side = max(bbox_width, bbox_height) * pad_factor
    half_side = square_side / 2
    
    # Calculate square crop bounds (can extend beyond image)
    crop_x0 = int(center_x - half_side)
    crop_y0 = int(center_y - half_side)
    crop_x1 = int(center_x + half_side)
    crop_y1 = int(center_y + half_side)
    
    # Get image dimensions
    img_width, img_height = image.size
    
    # Convert image to RGBA to support transparency
    if image.mode != 'RGBA':
        image = image.convert('RGBA')
    
    # Create a new square image with transparent background
    output_size = int(square_side)
    output_img = Image.new('RGBA', (output_size, output_size), (0, 0, 0, 0))
    
    # Calculate paste position (offset if crop extends beyond image bounds)
    paste_x = max(0, -crop_x0)
    paste_y = max(0, -crop_y0)
    
    # Calculate source crop region (clamped to image bounds)
    src_x0 = max(0, crop_x0)
    src_y0 = max(0, crop_y0)
    src_x1 = min(img_width, crop_x1)
    src_y1 = min(img_height, crop_y1)
    
    # Crop the region from source image
    if src_x1 > src_x0 and src_y1 > src_y0:
        cropped_region = image.crop((src_x0, src_y0, src_x1, src_y1))
        output_img.paste(cropped_region, (paste_x, paste_y))
    
    return output_img


class ImagingToolsUI:
    """Imaging utilities for the Rootweiler 'Imaging' section."""

    @classmethod
    def render(cls):
        st.subheader("Imaging")

        tabs = st.tabs(["Image extractor", "Compressor", "Object Crop (Square PNG)"])

        with tabs[0]:
            cls._render_extractor()

        with tabs[1]:
            cls._render_compressor()

        with tabs[2]:
            cls._render_object_crop()

    # ------------------------------------------------------------------
    # IMAGE EXTRACTOR
    # ------------------------------------------------------------------
    @classmethod
    def _render_extractor(cls):
        st.markdown(
            """
            ### Image extractor
            
            Upload a report or slide deck and pick which images you want to keep.
            Supports **PDF, PowerPoint (PPTX), Word (DOCX), and Excel (XLSX)**.
            """
        )

        uploaded = st.file_uploader(
            "Upload a single document",
            type=["pdf", "pptx", "docx", "xlsx"],
            key="imaging_extractor_uploader",
        )

        if uploaded is not None and st.button("Extract images", type="primary", key="extract_btn"):
            file_bytes = uploaded.getvalue()
            filename = uploaded.name
            try:
                images = cls._extract_images_from_file(file_bytes, filename)
            except Exception as e:
                st.error(f"Could not extract images: {e}")
                return

            if not images:
                st.warning("No images found (or all images were too small).")
                st.session_state["extracted_images"] = []
            else:
                # Store in session so selection persists across reruns
                st.session_state["extracted_images"] = [
                    {
                        "name": name,
                        "bytes": cls._pil_to_bytes(img),
                        "selected": True,
                    }
                    for img, name in images
                ]

        img_state = st.session_state.get("extracted_images", [])

        if not img_state:
            st.info("Extracted images will appear here after you upload and run extraction.")
            return

        st.markdown("#### Review and select images")

        # Select all / none
        c1, c2 = st.columns(2)
        with c1:
            if st.button("Select all", key="select_all_imgs"):
                for item in img_state:
                    item["selected"] = True
        with c2:
            if st.button("Deselect all", key="deselect_all_imgs"):
                for item in img_state:
                    item["selected"] = False

        # Display images in a grid with checkboxes
        cols_per_row = 3
        for i in range(0, len(img_state), cols_per_row):
            row_items = img_state[i : i + cols_per_row]
            cols = st.columns(len(row_items))
            for col, item, idx in zip(cols, row_items, range(i, i + len(row_items))):
                with col:
                    st.image(item["bytes"], use_column_width=True)
                    checked = st.checkbox(
                        item["name"],
                        value=item.get("selected", True),
                        key=f"img_select_{idx}",
                    )
                    item["selected"] = checked

        # Update state
        st.session_state["extracted_images"] = img_state

        # Prepare ZIP download of selected images
        selected = [item for item in img_state if item["selected"]]
        if not selected:
            st.warning("No images selected for download.")
            return

        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for item in selected:
                zf.writestr(item["name"], item["bytes"])
        zip_buffer.seek(0)

        st.download_button(
            "Download selected images (.zip)",
            data=zip_buffer,
            file_name="rootweiler_extracted_images.zip",
            mime="application/zip",
            key="download_extracted_zip",
        )

    # ----------------- extraction helpers ----------------- #

    @staticmethod
    def _extract_images_from_file(file_bytes: bytes, filename: str) -> List[Tuple[Image.Image, str]]:
        """Return list of (PIL.Image, suggested_name) from PDF / PPTX / DOCX / XLSX."""
        ext = os.path.splitext(filename)[1].lower()
        buf = io.BytesIO(file_bytes)

        images: List[Tuple[Image.Image, str]] = []

        if ext == ".pdf":
            doc = fitz.open(stream=buf.read(), filetype="pdf")
            for page_idx in range(len(doc)):
                page = doc.load_page(page_idx)
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    img_bytes = base_image["image"]
                    im = Image.open(io.BytesIO(img_bytes))
                    if im.width >= 30 and im.height >= 30:
                        name = f"pdf_page_{page_idx+1}_img_{img_index+1}.png"
                        images.append((im, name))

        elif ext == ".docx":
            doc = Document(buf)
            img_index = 0
            for rel in doc.part.rels:
                if "image" in doc.part.rels[rel].target_ref:
                    img_data = doc.part.rels[rel].target_part.blob
                    im = Image.open(io.BytesIO(img_data))
                    if im.width >= 30 and im.height >= 30:
                        img_index += 1
                        name = f"word_img_{img_index}.png"
                        images.append((im, name))

        elif ext == ".pptx":
            prs = Presentation(buf)
            img_index = 0
            for slide_idx, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_data = shape.image.blob
                        im = Image.open(io.BytesIO(img_data))
                        if im.width >= 30 and im.height >= 30:
                            img_index += 1
                            name = f"ppt_slide_{slide_idx}_img_{img_index}.png"
                            images.append((im, name))

        elif ext == ".xlsx":
            wb = openpyxl.load_workbook(buf)
            img_index = 0
            for sheet in wb:
                # _images is semi-private but commonly used
                for img in getattr(sheet, "_images", []):
                    try:
                        img_data = img.ref.blob  # may differ by openpyxl version
                    except AttributeError:
                        continue
                    im = Image.open(io.BytesIO(img_data))
                    if im.width >= 30 and im.height >= 30:
                        img_index += 1
                        safe_title = "".join(
                            c if c.isalnum() else "_" for c in sheet.title
                        )[:20]
                        name = f"xlsx_{safe_title}_img_{img_index}.png"
                        images.append((im, name))

        return images

    @staticmethod
    def _pil_to_bytes(img: Image.Image) -> bytes:
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()

    # ------------------------------------------------------------------
    # COMPRESSOR (same behavior as before)
    # ------------------------------------------------------------------
    @classmethod
    def _render_compressor(cls):
        st.markdown(
            """
            ### Image / file compressor
            
            Upload images or other files, choose a target size, and download a compressed ZIP.
            
            Images are compressed by:
            - Optional resizing to a maximum dimension
            - JPEG re-encoding with adjusted quality to aim for your target size
            
            Non-image files are included unchanged in the ZIP.
            """
        )

        uploaded_files = st.file_uploader(
            "Upload one or more files",
            type=None,
            accept_multiple_files=True,
            key="imaging_compressor_uploader",
        )

        if not uploaded_files:
            st.info("Add a few images or files above to get started.")
            return

        st.markdown("#### Compression settings")

        col1, col2 = st.columns(2)

        with col1:
            target_size_kb = st.slider(
                "Approximate target size per image (kB)",
                min_value=50,
                max_value=5000,
                value=500,
                step=50,
                help="The compressor will aim for this size per image by adjusting JPEG quality.",
            )

        with col2:
            resize_option = st.selectbox(
                "Max image dimension (longest side)",
                options=[
                    "No resizing",
                    "800 px",
                    "1200 px",
                    "2000 px",
                ],
                index=1,
                help="Resizing down first often gives better quality at smaller file sizes.",
            )

        max_dim = None
        if resize_option != "No resizing":
            max_dim = int(resize_option.split()[0])

        if st.button("Run compression", type="primary", key="compress_btn"):
            cls._run_compression(uploaded_files, target_size_kb, max_dim)

    @staticmethod
    def _run_compression(uploaded_files, target_size_kb: int, max_dim: int | None):
        compressed_files: list[tuple[str, bytes]] = []
        summary_rows = []

        for f in uploaded_files:
            raw_bytes = f.getvalue()
            orig_size_kb = len(raw_bytes) / 1024.0
            ext = os.path.splitext(f.name)[1].lower()

            is_image = False
            try:
                img = Image.open(io.BytesIO(raw_bytes))
                img.verify()
                is_image = True
            except Exception:
                is_image = False

            if is_image:
                img = Image.open(io.BytesIO(raw_bytes)).convert("RGB")

                compressed_bytes = ImagingToolsUI._compress_image(
                    img,
                    target_size_kb=target_size_kb,
                    max_dim=max_dim,
                )

                new_name = os.path.splitext(f.name)[0] + "_compressed.jpg"
                compressed_files.append((new_name, compressed_bytes))

                new_size_kb = len(compressed_bytes) / 1024.0
                ratio = new_size_kb / orig_size_kb if orig_size_kb > 0 else 1.0

                summary_rows.append(
                    {
                        "File": new_name,
                        "Original (kB)": round(orig_size_kb, 1),
                        "Compressed (kB)": round(new_size_kb, 1),
                        "Size ratio": f"{ratio:.2f}x",
                    }
                )
            else:
                compressed_files.append((f.name, raw_bytes))
                summary_rows.append(
                    {
                        "File": f.name,
                        "Original (kB)": round(orig_size_kb, 1),
                        "Compressed (kB)": round(orig_size_kb, 1),
                        "Size ratio": "1.00x (unchanged)",
                    }
                )

        if not compressed_files:
            st.warning("Nothing to compress – no valid files detected.")
            return

        st.markdown("#### Compression summary")
        st.dataframe(summary_rows, use_container_width=True)

        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for name, data in compressed_files:
                zf.writestr(name, data)
        zip_buffer.seek(0)

        st.download_button(
            "Download compressed files (.zip)",
            data=zip_buffer,
            file_name="rootweiler_compressed_files.zip",
            mime="application/zip",
            key="download_compressed_zip",
        )

    @staticmethod
    def _compress_image(
        img: Image.Image,
        target_size_kb: int,
        max_dim: int | None = None,
    ) -> bytes:
        """Resize (optional) and JPEG-compress an image, trying to hit target_size_kb."""
        if max_dim is not None:
            w, h = img.size
            longest = max(w, h)
            if longest > max_dim and longest > 0:
                scale = max_dim / float(longest)
                new_size = (int(w * scale), int(h * scale))
                img = img.resize(new_size, Image.LANCZOS)

        low_q, high_q = 20, 95
        best_bytes = None

        for _ in range(8):
            q = (low_q + high_q) // 2
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=q, optimize=True)
            data = buf.getvalue()
            size_kb = len(data) / 1024.0
            best_bytes = data

            if size_kb > target_size_kb and q > low_q:
                high_q = q - 1
            else:
                low_q = q + 1

        return best_bytes or b""

    # ------------------------------------------------------------------
    # OBJECT CROP (SQUARE PNG)
    # ------------------------------------------------------------------
    @classmethod
    def _render_object_crop(cls):
        st.markdown(
            """
            ### Object Crop (Square PNG)
            
            Upload an image, select an object region, and export a centered square PNG 
            with transparent padding if needed.
            """
        )

        uploaded_image = st.file_uploader(
            "Upload an image",
            type=["png", "jpg", "jpeg", "bmp", "webp"],
            key="object_crop_uploader",
        )

        if uploaded_image is None:
            st.info("Upload an image to get started.")
            return

        # Load the image
        try:
            img = Image.open(uploaded_image)
            img_array = np.array(img.convert('RGB'))
        except Exception as e:
            st.error(f"Could not load image: {e}")
            return

        st.markdown("#### Select object region")
        
        # Padding factor slider
        pad_factor = st.slider(
            "Padding factor",
            min_value=1.0,
            max_value=2.0,
            value=1.1,
            step=0.05,
            help="How much space around the object (1.0 = tight crop, 2.0 = 2x the object size)"
        )

        # Choose interaction method
        if CANVAS_AVAILABLE:
            st.caption("Draw a rectangle around the object you want to crop.")
            
            # Get image dimensions
            height, width = img_array.shape[:2]
            
            # Scale canvas to fit screen (max 700px wide)
            max_canvas_width = 700
            if width > max_canvas_width:
                scale = max_canvas_width / width
                canvas_width = max_canvas_width
                canvas_height = int(height * scale)
            else:
                scale = 1.0
                canvas_width = width
                canvas_height = height
            
            # Create canvas for drawing
            canvas_result = st_canvas(
                fill_color="rgba(255, 0, 0, 0.0)",  # Transparent fill
                stroke_width=2,
                stroke_color="#FF0000",
                background_image=Image.fromarray(img_array),
                update_streamlit=True,
                height=canvas_height,
                width=canvas_width,
                drawing_mode="rect",
                key="object_crop_canvas",
            )
            
            # Process canvas result
            if canvas_result.json_data is not None:
                objects = canvas_result.json_data.get("objects", [])
                
                if len(objects) == 0:
                    st.info("Draw a rectangle on the image to select the object region.")
                    return
                
                # Get the last drawn rectangle
                rect = objects[-1]
                
                # Extract coordinates and scale back to original image size
                canvas_x0 = rect["left"]
                canvas_y0 = rect["top"]
                canvas_width_rect = rect["width"]
                canvas_height_rect = rect["height"]
                
                # Scale back to original image coordinates
                x0 = int(canvas_x0 / scale)
                y0 = int(canvas_y0 / scale)
                x1 = int((canvas_x0 + canvas_width_rect) / scale)
                y1 = int((canvas_y0 + canvas_height_rect) / scale)
                
                bbox = (x0, y0, x1, y1)
                
                st.markdown(f"**Selected region:** ({x0}, {y0}) to ({x1}, {y1})")
                
            else:
                st.info("Draw a rectangle on the image to select the object region.")
                return
        else:
            # Fallback: manual bbox input
            st.warning("streamlit-drawable-canvas not available. Using manual input mode.")
            st.caption("Enter the bounding box coordinates manually.")
            
            col1, col2 = st.columns(2)
            with col1:
                x0 = st.number_input("Top-left X", min_value=0, max_value=img.width, value=0, step=1)
                y0 = st.number_input("Top-left Y", min_value=0, max_value=img.height, value=0, step=1)
            with col2:
                x1 = st.number_input("Bottom-right X", min_value=0, max_value=img.width, value=min(100, img.width), step=1)
                y1 = st.number_input("Bottom-right Y", min_value=0, max_value=img.height, value=min(100, img.height), step=1)
            
            bbox = (x0, y0, x1, y1)
            
            # Show preview with bbox
            preview_img = img.copy()
            draw = ImageDraw.Draw(preview_img)
            draw.rectangle([x0, y0, x1, y1], outline="red", width=3)
            st.image(preview_img, caption="Preview with bounding box", use_column_width=True)

        # Generate cropped square
        try:
            cropped_square = crop_square_png(img, bbox, pad_factor)
            
            st.markdown("#### Preview")
            st.image(cropped_square, caption=f"Cropped square ({cropped_square.width}×{cropped_square.height} px)", use_column_width=True)
            
            # Prepare download
            output_buffer = io.BytesIO()
            cropped_square.save(output_buffer, format="PNG")
            output_buffer.seek(0)
            
            # Generate filename
            original_name = os.path.splitext(uploaded_image.name)[0]
            output_filename = f"{original_name}_square_crop.png"
            
            st.download_button(
                "Download square PNG",
                data=output_buffer,
                file_name=output_filename,
                mime="image/png",
                key="download_square_crop",
            )
            
        except Exception as e:
            st.error(f"Could not generate cropped square: {e}")
