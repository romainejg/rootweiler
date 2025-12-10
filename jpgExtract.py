import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from tkinter.font import Font
from PIL import Image, ImageTk
import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import openpyxl
from io import BytesIO

class ImageExtractorApp:
    def __init__(self, root, back_callback):
        self.root = root
        self.back_callback = back_callback
        self.images = []
        self.current_image_index = 0
        self.current_image = None
        self.current_page_image = None
        self.rotated_image = None
        self.show_options()

    def show_options(self):
        self.clear_root_widgets()
        self.header_font = Font(family="Helvetica", size=28, weight="bold")
        self.subtext_font = Font(family="Helvetica", size=16, slant="italic")
        self.button_font = Font(family="Helvetica", size=16, weight="bold")

        frame = ttk.Frame(self.root, padding="20")
        frame.pack(fill=tk.BOTH, expand=True)

        header_label = ttk.Label(frame, text="Image Extractor", font=self.header_font)
        header_label.pack(pady=10)

        subtext_label = ttk.Label(frame, text="Select an option below", font=self.subtext_font)
        subtext_label.pack(pady=10)

        # Load and display logo image (e.g., a dog photo)
        script_dir = os.path.dirname(os.path.abspath(__file__))  # Get the directory of the script
        logo_path = os.path.join(script_dir, "logo.png")  # Construct the full path to the logo image

        try:
            print(f"Loading logo image from {logo_path}")
            logo_image = Image.open(logo_path)  # Load the logo image from the script directory
            logo_image = logo_image.resize((300, 300), Image.LANCZOS)
            logo_photo = ImageTk.PhotoImage(logo_image)
            logo_label = ttk.Label(frame, image=logo_photo)
            logo_label.image = logo_photo  # Keep a reference to avoid garbage collection
            logo_label.pack(pady=20)
        except Exception as e:
            print(f"Error loading logo image: {e}")

        style = ttk.Style()
        style.configure('TButton', font=self.button_font, padding=10)

        save_all_button = ttk.Button(frame, text="Save all img from file", command=self.save_all_images_prompt, style='TButton')
        save_all_button.pack(pady=10, fill=tk.X)

        selective_save_button = ttk.Button(frame, text="Choose which img to save from file", command=self.select_file, style='TButton')
        selective_save_button.pack(pady=10, fill=tk.X)

        back_button = ttk.Button(frame, text="Back", command=self.back_callback, style='TButton')
        back_button.pack(pady=10, fill=tk.X)




    def select_file(self):
        file_path = filedialog.askopenfilename(
            title="Select a File",
            filetypes=(
                ("All supported files", "*.pdf *.docx *.pptx *.xlsx"),
                ("PDF files", "*.pdf"),
                ("Word files", "*.docx"),
                ("PowerPoint files", "*.pptx"),
                ("Excel files", "*.xlsx")
            )
        )
        if file_path:
            self.process_file(file_path)

    def process_file(self, file_path):
        self.images = []  # Reset the images list
        self.current_image_index = 0  # Reset the image index
        extension = os.path.splitext(file_path)[1].lower()
        if extension == ".pdf":
            self.process_pdf(file_path)
        elif extension == ".docx":
            self.process_word(file_path)
        elif extension == ".pptx":
            self.process_pptx(file_path)
        elif extension == ".xlsx":
            self.process_xlsx(file_path)
        else:
            messagebox.showerror("Error", "Unsupported file type")

        if self.images:
            self.show_image(self.images[self.current_image_index])
        else:
            messagebox.showinfo("Info", "No images found in the selected file.")
            self.show_options()

    def process_pdf(self, file_path):
        doc = fitz.open(file_path)
        for i in range(len(doc)):
            page = doc.load_page(i)
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image = Image.open(BytesIO(image_bytes))
                if image.width >= 30 and image.height >= 30:  # Filter out small images
                    self.images.append((image, f"pdf_page_{i+1}_image_{img_index+1}.png", page))

    def process_word(self, file_path):
        doc = Document(file_path)
        image_index = 0
        for rel in doc.part.rels:
            if "image" in doc.part.rels[rel].target_ref:
                img_data = doc.part.rels[rel].target_part.blob
                image = Image.open(BytesIO(img_data))
                if image.width >= 30 and image.height >= 30:  # Filter out small images
                    image_index += 1
                    self.images.append((image, f"word_image_{image_index}.png", None))

    def process_pptx(self, file_path):
        prs = Presentation(file_path)
        image_index = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shape.image
                    img_data = img.blob
                    image = Image.open(BytesIO(img_data))
                    if image.width >= 30 and image.height >= 30:  # Filter out small images
                        image_index += 1
                        self.images.append((image, f"pptx_image_{image_index}.png", slide))

    def process_xlsx(self, file_path):
        wb = openpyxl.load_workbook(file_path)
        image_index = 0
        for sheet in wb:
            for image in sheet._images:
                img_data = image.ref.blob
                image = Image.open(BytesIO(img_data))
                if image.width >= 30 and image.height >= 30:  # Filter out small images
                    image_index += 1
                    self.images.append((image, f"xlsx_image_{sheet.title}_{image_index}.png", sheet))

    def show_image(self, image_data):
        self.clear_root_widgets()
        image, default_name, context = image_data
        self.current_image = image
        self.rotated_image = image  # Initialize rotated_image with the original image

        # Custom font for buttons
        button_font = Font(family="Helvetica", size=12, weight="bold")

        button_frame = ttk.Frame(self.root)
        button_frame.pack(pady=10)

        save_button = ttk.Button(button_frame, text="Save", command=lambda: self.save_image(self.rotated_image if self.rotated_image else image, default_name), style='TButton')
        save_button.grid(row=0, column=0, padx=10)

        skip_button = ttk.Button(button_frame, text="Skip", command=self.skip_image, style='TButton')
        skip_button.grid(row=0, column=1, padx=10)

        back_button = ttk.Button(button_frame, text="Back", command=self.back_image, style='TButton')
        back_button.grid(row=0, column=2, padx=10)

        rotate_left_button = ttk.Button(button_frame, text="Rotate Left", command=lambda: self.rotate_image(-90, right_frame, img_label), style='TButton')
        rotate_left_button.grid(row=0, column=3, padx=10)

        rotate_right_button = ttk.Button(button_frame, text="Rotate Right", command=lambda: self.rotate_image(90, right_frame, img_label), style='TButton')
        rotate_right_button.grid(row=0, column=4, padx=10)

        # Frame to hold context and image side by side
        content_frame = ttk.Frame(self.root)
        content_frame.pack(fill=tk.BOTH, expand=True)

        left_frame = ttk.Frame(content_frame, width=600, height=800)
        left_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=10, pady=10)
        left_frame.update_idletasks()
        left_frame.pack_propagate(False)

        right_frame = ttk.Frame(content_frame, width=600, height=800)
        right_frame.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True, padx=10, pady=10)
        right_frame.update_idletasks()
        right_frame.pack_propagate(False)

        if context:
            if isinstance(context, fitz.Page):
                # Display the entire page for PDF context
                page_image = context.get_pixmap()
                page_bytes = page_image.tobytes("png")
                page_image = Image.open(BytesIO(page_bytes))
            elif isinstance(context, Presentation):
                # Display the slide image for PowerPoint context
                slide_image = context.shapes[0].image
                slide_bytes = slide_image.blob
                page_image = Image.open(BytesIO(slide_bytes))
            else:
                # For other types, we do not have a visual representation
                page_image = None

            if page_image:
                self.current_page_image = page_image
                resized_page_image = self.fit_image_to_frame(page_image, left_frame)
                page_img = ImageTk.PhotoImage(resized_page_image)
                page_label = ttk.Label(left_frame, image=page_img)
                page_label.image = page_img
                page_label.pack(fill=tk.BOTH, expand=True)

        resized_image = self.fit_image_to_frame(image, right_frame)
        img = ImageTk.PhotoImage(resized_image)
        img_label = ttk.Label(right_frame, image=img)
        img_label.image = img
        img_label.pack(fill=tk.BOTH, expand=True)

        # Bind the resize event to dynamically resize images
        self.root.bind("<Configure>", lambda event: self.resize_images(event, left_frame, right_frame, page_label, img_label))

    def fit_image_to_frame(self, image, frame):
        frame.update_idletasks()
        frame_width = frame.winfo_width() or 600
        frame_height = frame.winfo_height() or 800

        # Calculate the aspect ratios
        image_ratio = image.width / image.height
        frame_ratio = frame_width / frame_height

        if image_ratio > frame_ratio:
            # Image is wider than the frame
            new_width = frame_width
            new_height = int(frame_width / image_ratio)
        else:
            # Image is taller than the frame
            new_height = frame_height
            new_width = int(frame_height * image_ratio)

        if new_width <= 0 or new_height <= 0:
            raise ValueError("Calculated image dimensions are invalid (width and height must be > 0)")

        return image.resize((new_width, new_height), Image.LANCZOS)

    def resize_images(self, event, left_frame, right_frame, page_label, img_label):
        if self.current_page_image and left_frame.winfo_width() > 0 and left_frame.winfo_height() > 0:
            resized_page_image = self.fit_image_to_frame(self.current_page_image, left_frame)
            page_img = ImageTk.PhotoImage(resized_page_image)
            page_label.configure(image=page_img)
            page_label.image = page_img

        if self.rotated_image and right_frame.winfo_width() > 0 and right_frame.winfo_height() > 0:
            resized_image = self.fit_image_to_frame(self.rotated_image, right_frame)
            img = ImageTk.PhotoImage(resized_image)
            img_label.configure(image=img)
            img_label.image = img

    def rotate_image(self, angle, frame, label):
        self.rotated_image = self.rotated_image.rotate(angle, expand=True)
        resized_image = self.fit_image_to_frame(self.rotated_image, frame)
        img = ImageTk.PhotoImage(resized_image)
        label.configure(image=img)
        label.image = img

    def save_image(self, image, default_name):
        save_path = filedialog.asksaveasfilename(defaultextension=".png", initialfile=default_name, filetypes=[("PNG files", "*.png")])
        if save_path:
            image.save(save_path)
            messagebox.showinfo("Success", f"Image saved as {save_path}")
        self.show_next_image()

    def skip_image(self):
        self.show_next_image()

    def back_image(self):
        if self.current_image_index > 0:
            self.current_image_index -= 1
            self.show_image(self.images[self.current_image_index])
        else:
            self.show_options()

    def show_next_image(self):
        self.current_image_index += 1
        if self.current_image_index < len(self.images):
            self.show_image(self.images[self.current_image_index])
        else:
            messagebox.showinfo("Info", "All images reviewed.")
            self.show_options()

    def save_all_images_prompt(self):
        file_path = filedialog.askopenfilename(
            title="Select a File",
            filetypes=(
                ("All supported files", "*.pdf *.docx *.pptx *.xlsx"),
                ("PDF files", "*.pdf"),
                ("Word files", "*.docx"),
                ("PowerPoint files", "*.pptx"),
                ("Excel files", "*.xlsx")
            )
        )
        if file_path:
            self.process_file_for_saving(file_path)
            self.save_all_images(file_path)

    def process_file_for_saving(self, file_path):
        self.images = []  # Reset the images list
        extension = os.path.splitext(file_path)[1].lower()
        if extension == ".pdf":
            self.process_pdf(file_path)
        elif extension == ".docx":
            self.process_word(file_path)
        elif extension == ".pptx":
            self.process_pptx(file_path)
        elif extension == ".xlsx":
            self.process_xlsx(file_path)
        else:
            messagebox.showerror("Error", "Unsupported file type")

    def save_all_images(self, file_path):
        save_folder = os.path.dirname(file_path)
        for img_data in self.images:
            image, default_name, _ = img_data
            save_path = os.path.join(save_folder, default_name)
            image.save(save_path)
        messagebox.showinfo("Success", f"All images saved to {save_folder}")
        self.show_options()

    def clear_root_widgets(self):
        for widget in self.root.winfo_children():
            widget.destroy()

def main_interface(root, back_callback):
    app = ImageExtractorApp(root, back_callback)
